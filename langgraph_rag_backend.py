"""
langgraph_rag_backend.py
Full RAG pipeline:
  Ingestion → Cleaning → Chunking → Embeddings → VectorDB
  → Hybrid Search → Query Transform → Re-ranking
  → Context Building → LLM → Observability → Evaluation
"""
from __future__ import annotations

import os, io, json, sqlite3, datetime, logging, time, re, unicodedata
from typing import Annotated, Any, Dict, Optional, TypedDict, List, Tuple

import yfinance as yf
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.tools.tavily_search import TavilySearchResults
from langchain_community.vectorstores import FAISS
from langchain_community.retrievers import BM25Retriever
from langchain_core.documents import Document
from langchain_core.messages import BaseMessage, SystemMessage, HumanMessage, AIMessage
from langchain_core.tools import tool
from langchain_huggingface import HuggingFaceEmbeddings
from langgraph.checkpoint.sqlite import SqliteSaver
from langgraph.graph import START, StateGraph
from langgraph.graph.message import add_messages
from langchain_groq import ChatGroq
from langgraph.prebuilt import ToolNode, tools_condition

load_dotenv()
logger = logging.getLogger(__name__)

from college_rag import build_college_index, college_rag_search, college_index_ready, get_college_index_meta

# ── 0. API key store ──────────────────────────────────────────────────────────
_KEY_DB_PATH = "user_api_keys.db"

def _get_key_conn() -> sqlite3.Connection:
    conn = sqlite3.connect(_KEY_DB_PATH, check_same_thread=False)
    conn.execute("""CREATE TABLE IF NOT EXISTS user_keys (
        user_id TEXT PRIMARY KEY, groq_key TEXT DEFAULT '', tavily_key TEXT DEFAULT '')""")
    conn.commit()
    return conn

_key_conn = _get_key_conn()

def save_user_keys(user_id: str, groq_key: str = "", tavily_key: str = "") -> None:
    _key_conn.execute(
        """INSERT INTO user_keys (user_id, groq_key, tavily_key) VALUES (?, ?, ?)
           ON CONFLICT(user_id) DO UPDATE SET groq_key=excluded.groq_key, tavily_key=excluded.tavily_key""",
        (user_id, groq_key.strip(), tavily_key.strip()))
    _key_conn.commit()

def load_user_keys(user_id: str) -> dict:
    row = _key_conn.execute("SELECT groq_key, tavily_key FROM user_keys WHERE user_id=?", (user_id,)).fetchone()
    return {"groq_key": row[0] or "", "tavily_key": row[1] or ""} if row else {"groq_key": "", "tavily_key": ""}

def user_has_keys(user_id: str) -> bool:
    return bool(load_user_keys(user_id)["groq_key"])

class MissingAPIKeyError(RuntimeError):
    pass

# ── 1. Embeddings + LLM factory ───────────────────────────────────────────────
_embeddings_instance = None
_college_index_attempted = False
_college_index_error = ""

def get_embeddings() -> HuggingFaceEmbeddings:
    global _embeddings_instance
    if _embeddings_instance is None:
        _embeddings_instance = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    return _embeddings_instance

def ensure_college_index() -> None:
    global _college_index_attempted, _college_index_error
    if _college_index_attempted:
        return
    _college_index_attempted = True
    try:
        build_college_index(get_embeddings())
        _college_index_error = ""
    except Exception as exc:
        _college_index_error = str(exc)
        logger.exception("College knowledge base index failed to build.")

def get_llm(user_id: str) -> ChatGroq:
    api_key = load_user_keys(user_id)["groq_key"] or os.getenv("GROQ_API_KEY", "")
    if not api_key:
        raise MissingAPIKeyError("No Groq API key found. Please click **API Keys**.")
    return ChatGroq(model="llama-3.3-70b-versatile", api_key=api_key)

def get_search_tool(user_id: str) -> TavilySearchResults:
    os.environ["TAVILY_API_KEY"] = load_user_keys(user_id)["tavily_key"] or os.getenv("TAVILY_API_KEY", "")
    return TavilySearchResults(max_results=3)

_default_llm_instance = None
def _get_default_llm() -> ChatGroq:
    global _default_llm_instance
    api_key = os.getenv("GROQ_API_KEY", "")
    if not api_key:
        raise MissingAPIKeyError("No Groq API key found. Please click **API Keys**.")
    if _default_llm_instance is None:
        _default_llm_instance = ChatGroq(model="llama-3.3-70b-versatile", api_key=api_key)
    return _default_llm_instance

def _llm_call(user_id: str, system_prompt: str, user_prompt: str) -> str:
    llm = get_llm(user_id) if user_id else _get_default_llm()
    return llm.invoke([SystemMessage(content=system_prompt), {"role": "user", "content": user_prompt}]).content.strip()

def search_college_kb(query: str) -> dict:
    ensure_college_index()
    if _college_index_error:
        return {
            "error": f"College knowledge base could not initialize: {_college_index_error}",
            "query": query,
            "context": [],
        }
    return college_rag_search(query)

# ── 2. INGESTION: Universal file → clean Markdown ────────────────────────────

# CHANGED: explicit set of image extensions we refuse to ingest as documents.
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"}

def _clean_text(text: str) -> str:
    """
    CLEANING stage:
    - Normalise unicode (NFC)
    - Remove null bytes and non-printable control characters
    - Collapse excessive whitespace / blank lines
    - Strip page-header/footer artifacts (lines that are purely numbers or short repeated tokens)
    """
    text = unicodedata.normalize("NFC", text)
    text = text.replace("\x00", "")
    text = re.sub(r"[^\x09\x0A\x0D\x20-\x7E\u00A0-\uFFFF]", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [l for l in text.splitlines()
             if not re.fullmatch(r"\s*[\d\-–—]+\s*", l)]
    return "\n".join(lines).strip()

def _pdf_to_markdown(file_bytes: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"## Page {i}\n\n{_clean_text(text)}")
    return "\n\n---\n\n".join(parts)

def _docx_to_markdown(file_bytes: bytes) -> str:
    from docx import Document as DocxDoc
    doc = DocxDoc(io.BytesIO(file_bytes))
    lines = []
    for para in doc.paragraphs:
        s = para.style.name.lower(); t = para.text.strip()
        if not t: lines.append(""); continue
        if "heading 1" in s: lines.append(f"# {t}")
        elif "heading 2" in s: lines.append(f"## {t}")
        elif "heading 3" in s: lines.append(f"### {t}")
        else: lines.append(t)
    for table in doc.tables:
        if not table.rows: continue
        hdr = [c.text.strip() for c in table.rows[0].cells]
        lines += ["\n| "+" | ".join(hdr)+" |", "| "+" | ".join(["---"]*len(hdr))+" |"]
        for row in table.rows[1:]:
            lines.append("| "+" | ".join(c.text.strip() for c in row.cells)+" |")
    return _clean_text("\n".join(lines))

def _pptx_to_markdown(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip(): lines.append(shape.text.strip())
        parts.append("\n\n".join(lines))
    return _clean_text("\n\n---\n\n".join(parts))

def _xlsx_to_markdown(file_bytes: bytes) -> str:
    import pandas as pd
    xl = pd.ExcelFile(io.BytesIO(file_bytes))
    return "\n\n---\n\n".join(f"## Sheet: {s}\n\n{xl.parse(s).to_markdown(index=False)}" for s in xl.sheet_names)

def _csv_to_markdown(file_bytes: bytes) -> str:
    import pandas as pd
    return pd.read_csv(io.BytesIO(file_bytes)).to_markdown(index=False)

def _txt_to_markdown(file_bytes: bytes) -> str:
    return _clean_text(file_bytes.decode("utf-8", errors="replace"))

def convert_file_to_markdown(file_bytes: bytes, filename: str) -> Tuple[str, str]:
    ext = os.path.splitext(filename.lower())[1]
    # CHANGED: reject images explicitly instead of mis-decoding them as text.
    if ext in IMAGE_EXTS:
        raise ValueError("Images can't be ingested as documents. Use the Diagram tab to explain images.")
    if ext == ".pdf": return _pdf_to_markdown(file_bytes), "pdf"
    elif ext == ".docx": return _docx_to_markdown(file_bytes), "docx"
    elif ext == ".pptx": return _pptx_to_markdown(file_bytes), "pptx"
    elif ext in (".xlsx", ".xls"): return _xlsx_to_markdown(file_bytes), "xlsx"
    elif ext == ".csv": return _csv_to_markdown(file_bytes), "csv"
    else: return _txt_to_markdown(file_bytes), "text"

# ── 3. CHUNKING ───────────────────────────────────────────────────────────────
def _make_splitter(chunk_size: int = 800, chunk_overlap: int = 150) -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n## ", "\n### ", "\n\n", "\n", ". ", " ", ""],
    )

# ── 4. Per-thread vector store (FAISS + BM25 for hybrid search) ───────────────
_THREAD_VS: Dict[str, FAISS]         = {}
_THREAD_BM25: Dict[str, BM25Retriever] = {}
_THREAD_CHUNKS: Dict[str, List[Document]] = {}
_THREAD_METADATA: Dict[str, dict]    = {}
_THREAD_TITLES: Dict[str, str]       = {}

# ── 5. OBSERVABILITY store ────────────────────────────────────────────────────
_OBS_LOG: List[dict] = []

def _obs_log(entry: dict):
    _OBS_LOG.append({**entry, "ts": datetime.datetime.utcnow().isoformat()})
    if len(_OBS_LOG) > 500: _OBS_LOG.pop(0)

def get_observability_log() -> List[dict]:
    return list(reversed(_OBS_LOG))

def _get_thread_vs(thread_id: Optional[str]) -> Optional[FAISS]:
    return _THREAD_VS.get(str(thread_id)) if thread_id else None

def ingest_file(file_bytes: bytes, thread_id: str, filename: Optional[str] = None) -> dict:
    """Full pipeline: Ingest → Clean → Chunk → Embed → Store (FAISS + BM25)."""
    if not file_bytes: raise ValueError("No bytes received.")
    fname = filename or "uploaded_file"
    t0 = time.time()

    markdown_text, file_type = convert_file_to_markdown(file_bytes, fname)
    if not markdown_text.strip(): raise ValueError(f"Could not extract text from `{fname}`.")

    splitter = _make_splitter()
    chunks = splitter.split_documents([Document(page_content=markdown_text,
        metadata={"source": fname, "file_type": file_type, "thread_id": thread_id})])

    vs = FAISS.from_documents(chunks, get_embeddings())
    _THREAD_VS[str(thread_id)] = vs

    bm25 = BM25Retriever.from_documents(chunks)
    bm25.k = 6
    _THREAD_BM25[str(thread_id)] = bm25
    _THREAD_CHUNKS[str(thread_id)] = chunks

    elapsed = round((time.time() - t0) * 1000)
    meta = {"filename": fname, "file_type": file_type,
            "chunks": len(chunks), "char_count": len(markdown_text), "ingest_ms": elapsed}
    _THREAD_METADATA[str(thread_id)] = meta
    logger.info("ingest_file: %s → %d chunks in %dms", fname, len(chunks), elapsed)
    return meta

def ingest_pdf(file_bytes: bytes, thread_id: str, filename: Optional[str] = None) -> dict:
    return ingest_file(file_bytes, thread_id, filename)

# ── 6. HYBRID SEARCH + QUERY TRANSFORM + RE-RANKING ──────────────────────────

def _query_transform(query: str, llm: ChatGroq) -> List[str]:
    try:
        resp = llm.invoke([
            SystemMessage(content=(
                "You are a query expansion assistant. Given a question, produce 2 alternative "
                "phrasings that cover the same information need. "
                "Reply ONLY with a JSON array of 2 strings, no explanation."
            )),
            {"role": "user", "content": query},
        ])
        raw = resp.content.strip().replace("```json","").replace("```","").strip()
        alts = json.loads(raw)
        if isinstance(alts, list) and len(alts) >= 2:
            return [query] + list(alts[:2])
    except Exception:
        pass
    return [query]


def _reciprocal_rank_fusion(ranked_lists: List[List[Document]], k: int = 60) -> List[Document]:
    scores: Dict[str, float] = {}
    doc_map: Dict[str, Document] = {}
    for ranked in ranked_lists:
        for rank, doc in enumerate(ranked):
            key = doc.page_content[:200]
            scores[key] = scores.get(key, 0.0) + 1.0 / (k + rank + 1)
            doc_map[key] = doc
    sorted_keys = sorted(scores, key=scores.__getitem__, reverse=True)
    return [doc_map[k] for k in sorted_keys]


def _hybrid_search(thread_id: str, query: str, llm: Optional[ChatGroq] = None,
                   top_k: int = 5) -> Tuple[List[Document], dict]:
    vs   = _THREAD_VS.get(str(thread_id))
    bm25 = _THREAD_BM25.get(str(thread_id))
    if vs is None:
        return [], {}

    t0 = time.time()
    queries = _query_transform(query, llm) if llm else [query]

    all_lists: List[List[Document]] = []
    for q in queries:
        try:
            all_lists.append(vs.similarity_search(q, k=6))
        except Exception: pass
        if bm25:
            try:
                all_lists.append(bm25.invoke(q))
            except Exception: pass

    retrieval_ms = round((time.time() - t0) * 1000)
    t1 = time.time()
    reranked = _reciprocal_rank_fusion(all_lists) if all_lists else []
    final = reranked[:top_k]
    rerank_ms = round((time.time() - t1) * 1000)

    obs = {"queries_used": queries, "retrieval_ms": retrieval_ms,
           "rerank_ms": rerank_ms, "chunks_returned": len(final)}
    return final, obs


def hybrid_context(thread_id: str, query: str, llm: Optional[ChatGroq] = None) -> dict:
    docs, obs = _hybrid_search(thread_id, query, llm)
    _obs_log({**obs, "thread_id": thread_id, "query": query, "source": "user_doc"})
    if not docs:
        return {"error": "No document uploaded yet.", "query": query}
    return {"query": query, "context": [d.page_content for d in docs],
            "metadata": [d.metadata for d in docs], "obs": obs,
            "source_file": _THREAD_METADATA.get(str(thread_id), {}).get("filename")}

# ── 7. Thread title helpers ───────────────────────────────────────────────────
def set_thread_title(thread_id: str, title: str): _THREAD_TITLES[thread_id] = title
def get_thread_title(thread_id: str) -> str: return _THREAD_TITLES.get(thread_id, "")

def generate_chat_title(first_message: str, user_id: str = "") -> str:
    try:
        llm = get_llm(user_id) if user_id else _get_default_llm()
        resp = llm.invoke([
            SystemMessage(content="You are a chat title generator. Given the user's first message, respond with ONLY a short title of 3–5 words. No punctuation, no quotes, no explanation."),
            {"role": "user", "content": first_message},
        ])
        words = resp.content.strip().strip('"\'').split()
        return " ".join(words[:6]) if words else "New Chat"
    except Exception:
        return first_message.strip()[:40] + ("…" if len(first_message) > 40 else "")

# ── 8. LangGraph tools ────────────────────────────────────────────────────────
@tool
def calculator(first_num: float, second_num: float, operation: str) -> dict:
    """Perform basic arithmetic. Supported: add, sub, mul, div"""
    try:
        if operation == "add": r = first_num + second_num
        elif operation == "sub": r = first_num - second_num
        elif operation == "mul": r = first_num * second_num
        elif operation == "div":
            if second_num == 0: return {"error": "Division by zero"}
            r = first_num / second_num
        else: return {"error": f"Unsupported: {operation}"}
        return {"first_num": first_num, "second_num": second_num, "operation": operation, "result": r}
    except Exception as e: return {"error": str(e)}

@tool
def get_stock_price(symbol: str) -> dict:
    """Fetch latest stock price using Yahoo Finance."""
    ticker = yf.Ticker(symbol)
    info = ticker.fast_info
    return {"symbol": symbol, "current_price": info.last_price, "currency": info.currency}

# REMOVED: standalone @tool rag_tool that required the model to pass thread_id.
# It is now created per-thread inside get_chatbot via _make_rag_tool() so the
# thread_id is injected from server config and can never be omitted by the LLM.

@tool
def college_rag_tool(query: str) -> dict:
    """Search the college knowledge base (all semesters, lecture notes, syllabus).
    Use for ANY question about subjects, topics, exams, or coursework."""
    t0 = time.time()
    result = search_college_kb(query)
    _obs_log({"query": query, "source": "college_kb", "retrieval_ms": round((time.time()-t0)*1000),
              "chunks_returned": len(result.get("context", []))})
    return result

# CHANGED: factory that binds thread_id into the rag_tool closure.
def _make_rag_tool(thread_id: Optional[str]):
    @tool
    def rag_tool(query: str) -> dict:
        """Retrieve information from the user's uploaded document via hybrid search.
        The document is selected automatically for the current conversation."""
        return hybrid_context(str(thread_id) if thread_id else "", query)
    return rag_tool

# ── 9. Academic feature functions ─────────────────────────────────────────────

def generate_cheat_sheet(subject: str, unit: str, user_id: str = "") -> str:
    kb = search_college_kb(f"{subject} {unit} formulas definitions key concepts")
    context = "\n\n".join(kb.get("context", [])[:5])
    return _llm_call(user_id,
        "You are a concise academic summariser. Create compact cheat sheets using markdown "
        "with tables, bullet points, and formulas. Prioritise high-yield exam content.",
        f"Create an exam cheat sheet for: **{subject}** — {unit}\n\nSource:\n{context}\n\n"
        "Include: key definitions, important formulas, common exam questions, memory tricks.")

def detect_subject_unit(message: str, user_id: str = "") -> dict:
    try:
        raw = _llm_call(user_id,
            "You are an academic topic classifier for an Indian engineering college. "
            "Given a student's question, identify the most likely semester (1-8), subject, and unit. "
            'Reply ONLY with valid JSON: {"semester":3,"subject":"Data Structures","unit":"Trees","confidence":0.9} '
            "Use null for semester if unsure.",
            message).replace("```json","").replace("```","").strip()
        return json.loads(raw)
    except Exception:
        return {"semester": None, "subject": "Unknown", "unit": "Unknown", "confidence": 0.0}

def generate_flashcards(subject: str, unit: str, count: int = 10, user_id: str = "") -> List[dict]:
    """Generate Q&A flashcard pairs from KB. Returns list of {question, answer}."""
    kb = search_college_kb(f"{subject} {unit}")
    context = "\n\n".join(kb.get("context", [])[:6])
    raw = _llm_call(user_id,
        "You are a flashcard generator for college students. "
        f"Generate exactly {count} question-answer flashcard pairs from the provided notes. "
        "Reply ONLY with a valid JSON array: "
        '[{"question":"What is X?","answer":"X is..."},...]  No extra text.',
        f"Subject: {subject} | Unit: {unit}\n\nNotes:\n{context}")
    try:
        raw = raw.replace("```json","").replace("```","").strip()
        cards = json.loads(raw)
        if isinstance(cards, list): return cards[:count]
    except Exception:
        pass
    return [{"question": "Could not generate flashcards.", "answer": "Check your API key and KB content."}]

def explain_concept(concept: str, level: str = "intermediate", user_id: str = "") -> str:
    kb = search_college_kb(concept)
    context = "\n\n".join(kb.get("context", [])[:4])
    level_map = {
        "beginner":      "Explain like the student has zero background. Use simple analogies, avoid jargon, use real-world examples.",
        "intermediate":  "Explain clearly for a student who knows the basics. Use standard terminology, step-by-step reasoning.",
        "exam-ready":    "Explain concisely for exam revision: definition, key points, formula if any, 1 example, common exam traps.",
    }
    instruction = level_map.get(level, level_map["intermediate"])
    return _llm_call(user_id,
        f"You are an expert academic tutor. {instruction} Format with markdown.",
        f"Concept: **{concept}**\n\nRelevant notes:\n{context}")

def summarise_topic(topic: str, length: str = "paragraph", user_id: str = "") -> str:
    kb = search_college_kb(topic)
    context = "\n\n".join(kb.get("context", [])[:5])
    length_map = {
        "sentence":  "Summarise in exactly 1 sentence (max 30 words).",
        "paragraph": "Summarise in 1 clear paragraph (4–6 sentences).",
        "full":      "Write a full, structured summary using markdown headings and bullet points.",
    }
    instruction = length_map.get(length, length_map["paragraph"])
    return _llm_call(user_id,
        f"You are an academic summariser. {instruction}",
        f"Topic: **{topic}**\n\nSource material:\n{context}")

def explain_diagram_image(image_bytes: bytes, mime_type: str, user_id: str = "") -> str:
    import base64
    api_key = load_user_keys(user_id)["groq_key"] or os.getenv("GROQ_API_KEY", "")
    if not api_key:
        raise MissingAPIKeyError("No Groq API key. Please click **API Keys**.")

    vision_llm = ChatGroq(model="meta-llama/llama-4-scout-17b-16e-instruct", api_key=api_key)
    b64 = base64.b64encode(image_bytes).decode()
    msg = HumanMessage(content=[
        {"type": "text", "text": (
            "You are an academic diagram explainer. Analyse this diagram/figure/chart/whiteboard image and provide:\n"
            "1. **What it shows** (1 sentence)\n"
            "2. **Key components** (labelled list)\n"
            "3. **Step-by-step explanation** of the process or concept shown\n"
            "4. **Subject area** this belongs to\n"
            "5. **Exam tip** — what questions this diagram might appear in\n\n"
            "Use clear markdown formatting."
        )},
        {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{b64}"}},
    ])
    try:
        response = vision_llm.invoke([msg])
        return response.content.strip()
    except Exception as e:
        return f"Could not analyse image: {e}"

# REMOVED: export_conversation_pdf() and the fpdf dependency (PDF export dropped).

def export_conversation_txt(messages: list, thread_title: str = "Chat Export") -> str:
    lines = [thread_title, "="*len(thread_title),
             f"Exported: {datetime.datetime.now().strftime('%d %B %Y, %I:%M %p')}", "", ""]
    for msg in messages:
        role = "YOU" if msg.get("role") == "user" else "ASSISTANT"
        lines += [f"[{role}]", msg.get("content",""), "", "-"*60, ""]
    return "\n".join(lines)

# ── 10. EVALUATION helper ─────────────────────────────────────────────────────
def evaluate_retrieval(query: str, retrieved_docs: List[Document],
                       user_id: str = "") -> dict:
    if not retrieved_docs:
        return {"relevance_score": 0.0, "assessment": "No documents retrieved."}
    snippets = "\n---\n".join(d.page_content[:300] for d in retrieved_docs[:3])
    try:
        raw = _llm_call(user_id,
            "You are a retrieval evaluator. Given a query and retrieved chunks, "
            "score relevance from 0.0 to 1.0. Reply ONLY with JSON: "
            '{"relevance_score": 0.85, "assessment": "Chunks are highly relevant..."}',
            f"Query: {query}\n\nChunks:\n{snippets}")
        raw = raw.replace("```json","").replace("```","").strip()
        return json.loads(raw)
    except Exception:
        return {"relevance_score": 0.5, "assessment": "Could not evaluate."}

# ── 11. Graph builder ─────────────────────────────────────────────────────────
_COMPILED_GRAPHS: Dict[str, Any] = {}
conn = sqlite3.connect(":memory:", check_same_thread=False)
checkpointer = SqliteSaver(conn=conn)

def get_chatbot(user_id: str):
    if user_id in _COMPILED_GRAPHS:
        return _COMPILED_GRAPHS[user_id]

    user_llm     = get_llm(user_id)
    user_search  = get_search_tool(user_id)

    # CHANGED: rag_tool is built once per compiled graph. thread_id is resolved
    # at call time from config inside chat_node, so we bind a thin wrapper that
    # reads it dynamically. We register a placeholder tool and a ToolNode that
    # rebuilds with the live thread_id.
    class ChatState(TypedDict):
        messages: Annotated[List[BaseMessage], add_messages]

    def chat_node(state: ChatState, config=None):
        thread_id = (config or {}).get("configurable", {}).get("thread_id") if config else None

        # CHANGED: construct tools with thread_id injected from config.
        rag_tool = _make_rag_tool(thread_id)
        turn_tools = [get_stock_price, user_search, calculator, rag_tool, college_rag_tool]
        llm_tools = user_llm.bind_tools(turn_tools)

        kb_status = ("The college KB is loaded and ready." if is_college_kb_ready()
                     else "The college KB has no data yet.")
        sys_msg = SystemMessage(content=(
            "You are a helpful, friendly academic assistant for college students — "
            "respond conversationally like ChatGPT would, but grounded in course notes.\n\n"
            f"{kb_status}\n\n"
            "RULES:\n"
            "• For ANY academic question → call `college_rag_tool` first.\n"
            # CHANGED: no longer tell the model to pass thread_id; it's automatic.
            "• For the user's uploaded file → call `rag_tool` (the current document is selected automatically).\n"
            "• For live web info → use the search tool.\n"
            "• For math → use calculator.\n"
            "• After getting context, write a thorough, well-structured answer.\n"
            "• End with: 'Confidence: High/Medium/Low — <brief reason>.'"
        ))
        response = llm_tools.invoke([sys_msg, *state["messages"]], config=config)
        return {"messages": [response]}

    # CHANGED: ToolNode needs the same tool set. thread_id inside rag_tool is
    # resolved per-invocation through config, so we build a ToolNode whose
    # rag_tool reads the thread_id from the runnable config at call time.
    def tools_node(state: ChatState, config=None):
        thread_id = (config or {}).get("configurable", {}).get("thread_id") if config else None
        rag_tool = _make_rag_tool(thread_id)
        node = ToolNode([get_stock_price, user_search, calculator, rag_tool, college_rag_tool])
        return node.invoke(state, config=config)

    graph = StateGraph(ChatState)
    graph.add_node("chat_node", chat_node)
    graph.add_node("tools", tools_node)   # CHANGED: dynamic tools node
    graph.add_edge(START, "chat_node")
    graph.add_conditional_edges("chat_node", tools_condition)
    graph.add_edge("tools", "chat_node")

    compiled = graph.compile(checkpointer=checkpointer)
    _COMPILED_GRAPHS[user_id] = compiled
    return compiled

def invalidate_graph(user_id: str):
    _COMPILED_GRAPHS.pop(user_id, None)

chatbot = None  # legacy alias

# ── 12. Thread utilities ──────────────────────────────────────────────────────
def retrieve_user_threads(user_id: str) -> list:
    threads = []
    for cp in checkpointer.list(None):
        tid = cp.config["configurable"]["thread_id"]
        if tid.startswith(f"{user_id}::") and tid not in threads:
            threads.append(tid)
    return threads

def thread_has_document(thread_id: str) -> bool:
    return str(thread_id) in _THREAD_VS

def thread_document_metadata(thread_id: str) -> dict:
    return _THREAD_METADATA.get(str(thread_id), {})

def get_college_kb_meta() -> dict:
    ensure_college_index()
    if _college_index_error:
        return {"files": [], "chunks": 0, "chars": 0, "error": _college_index_error}
    return get_college_index_meta()

def is_college_kb_ready() -> bool:
    ensure_college_index()
    return college_index_ready()
